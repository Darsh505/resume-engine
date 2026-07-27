#!/usr/bin/env python3
"""
generate_presentation.py
────────────────────────
Generates a polished, 16-slide PowerPoint presentation for the
Resume Engine Summer Training Project.

Run:
    python generate_presentation.py
Output:
    Resume_Engine_Presentation.pptx
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ─── Color Palette ────────────────────────────────────────────────────────────
NAVY_BG        = RGBColor(0x0D, 0x11, 0x17)   # #0D1117  – deep dark background
NAVY_CARD      = RGBColor(0x16, 0x1B, 0x22)   # #161B22  – card surface
NAVY_MID       = RGBColor(0x1A, 0x1F, 0x2E)   # #1A1F2E  – mid tone
ACCENT_BLUE    = RGBColor(0x25, 0x63, 0xEB)   # #2563EB  – primary accent
ACCENT_CYAN    = RGBColor(0x38, 0xBD, 0xF8)   # #38BDF8  – secondary accent
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE      = RGBColor(0xE2, 0xE8, 0xF0)   # #E2E8F0  – body text
MUTED          = RGBColor(0x8B, 0x94, 0xA3)   # muted gray
GREEN_OK       = RGBColor(0x22, 0xC5, 0x5E)   # #22C55E
AMBER          = RGBColor(0xF5, 0x9E, 0x0B)   # #F59E0B

# ─── Slide dimensions (16:9 widescreen) ──────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ─── Typography Sizes ─────────────────────────────────────────────────────────
SZ_HERO        = Pt(52)
SZ_TITLE       = Pt(36)
SZ_SUBTITLE    = Pt(22)
SZ_SECTION     = Pt(18)
SZ_BODY        = Pt(14)
SZ_SMALL       = Pt(11)
SZ_CAPTION     = Pt(9)


# ─── Utility helpers ──────────────────────────────────────────────────────────

def _rgb(r, g, b): return RGBColor(r, g, b)

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None, radius=None):
    """Add a rectangle (optionally rounded) with fill/line settings."""
    from pptx.util import Emu as _Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()

    if radius:
        # Apply rounded corners via XML patch
        sp = shape._element
        sp_pr = sp.find(qn('p:spPr'))
        prstGeom = sp_pr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            # Update existing geometry to roundRect
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            # Remove old gd elements and add new one
            for gd in avLst.findall(qn('a:gd')):
                avLst.remove(gd)
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius}')
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_name="Aptos", font_size=Pt(14), bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT, wrap=True, line_spacing=None):
    """Add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if line_spacing:
        from pptx.util import Pt as Pt2
        from pptx.oxml.ns import qn
        from lxml import etree
        # Set line spacing
        pPr = p._pPr
        if pPr is None:
            pPr = etree.SubElement(p._p, qn('a:pPr'))
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 1000)))
    return txBox


def add_paragraph(tf, text, font_name="Aptos", font_size=Pt(13), bold=False,
                  italic=False, color=OFF_WHITE, align=PP_ALIGN.LEFT, space_before=0):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def set_slide_bg(slide, color=NAVY_BG):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_header(slide, title_text, accent=ACCENT_CYAN, number=None):
    """Add the standard slide title bar with accent underline."""
    # Background is already set; add top accent bar
    bar = add_shape(slide, 0, 0, SLIDE_W, Inches(0.07), fill_color=accent)

    # Slide number chip
    if number:
        chip = add_shape(slide, Inches(12.5), Inches(0.15), Inches(0.6), Inches(0.35),
                         fill_color=NAVY_CARD, radius=20000)
        add_textbox(slide, str(number),
                    Inches(12.5), Inches(0.15), Inches(0.6), Inches(0.35),
                    font_size=Pt(9), color=MUTED, align=PP_ALIGN.CENTER)

    # Title text
    add_textbox(slide, title_text,
                Inches(0.5), Inches(0.18), Inches(11.5), Inches(0.65),
                font_name="Aptos Display", font_size=Pt(28), bold=True,
                color=WHITE, align=PP_ALIGN.LEFT)

    # Accent underline
    add_shape(slide, Inches(0.5), Inches(0.82), Inches(2.5), Inches(0.04),
              fill_color=accent)


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    add_shape(slide, left, top, width, Inches(0.03), fill_color=color)


def add_card(slide, left, top, width, height, fill=NAVY_CARD, border=None, radius=30000):
    """Add a rounded card."""
    return add_shape(slide, left, top, width, height,
                     fill_color=fill, line_color=border, radius=radius)


def bullet_paragraph(tf, text, indent=0, font_size=Pt(12.5), color=OFF_WHITE, bullet_char="▸"):
    """Add a bullet point paragraph."""
    p = tf.add_paragraph()
    p.level = indent
    run = p.add_run()
    run.text = f"  {bullet_char}  {text}"
    run.font.name = "Aptos"
    run.font.size = font_size
    run.font.color.rgb = color
    return p


def pill_label(slide, text, left, top, width=Inches(1.8), height=Inches(0.32),
               bg=ACCENT_BLUE, fg=WHITE, font_size=Pt(10.5)):
    """Add a pill/badge label."""
    add_card(slide, left, top, width, height, fill=bg, radius=40000)
    add_textbox(slide, text, left, top, width, height,
                font_size=font_size, color=fg, align=PP_ALIGN.CENTER, bold=True)


def section_chip(slide, text, left, top):
    """Small section header chip."""
    add_shape(slide, left, top, Inches(2.0), Inches(0.28),
              fill_color=_rgb(0x25, 0x63, 0xEB), radius=15000)
    add_textbox(slide, text, left, top, Inches(2.0), Inches(0.28),
                font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def flow_arrow(slide, x, y, label=None, direction="right"):
    """Draw a compact flow arrow."""
    # Arrow body
    add_shape(slide, x, y + Inches(0.085), Inches(0.35), Inches(0.05),
              fill_color=ACCENT_CYAN)
    # Arrowhead (triangle approximation via text)
    add_textbox(slide, "▶", x + Inches(0.25), y, Inches(0.2), Inches(0.25),
                font_size=Pt(10), color=ACCENT_CYAN, align=PP_ALIGN.LEFT)


# ─── Speaker notes helper ─────────────────────────────────────────────────────

def add_speaker_notes(slide, notes_text: str):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


# ─── Individual Slide Builders ───────────────────────────────────────────────

def build_slide_01_title(prs):
    """Slide 1: Title / Hero"""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY_BG)

    # Top accent gradient bar
    add_shape(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=ACCENT_BLUE)

    # Left decorative vertical bar
    add_shape(slide, 0, 0, Inches(0.06), SLIDE_H, fill_color=ACCENT_BLUE)

    # Abstract geometric decoration — top-right corner cluster
    # Large circle (gear/document motif)
    from pptx.util import Inches as I
    add_shape(slide, I(10.8), I(-0.5), I(3.2), I(3.2),
              fill_color=_rgb(0x1E, 0x3A, 0x5F), radius=100000)
    add_shape(slide, I(11.2), I(-0.1), I(2.4), I(2.4),
              fill_color=_rgb(0x25, 0x63, 0xEB), radius=100000)
    add_shape(slide, I(11.6), I(0.3), I(1.6), I(1.6),
              fill_color=_rgb(0x38, 0xBD, 0xF8), radius=100000)

    # Tag chips (decorative, floating)
    for i, (tag, tx, ty) in enumerate([
        ("#backend",     I(9.6),  I(1.2)),
        ("#frontend",    I(10.9), I(2.0)),
        ("#data-science",I(9.0),  I(2.2)),
        ("#devops",      I(11.2), I(3.0)),
        ("general",      I(9.8),  I(3.2)),
    ]):
        alpha = [ACCENT_BLUE, ACCENT_CYAN, _rgb(0x6D,0x28,0xD9),
                 _rgb(0x10,0xB9,0x81), _rgb(0xF5,0x9E,0x0B)][i]
        add_card(slide, tx, ty, I(1.6), I(0.3), fill=alpha, radius=30000)
        add_textbox(slide, tag, tx, ty, I(1.6), I(0.3),
                    font_size=Pt(9), color=WHITE, align=PP_ALIGN.CENTER, bold=True)

    # Bottom-left decorative dots
    for col in range(6):
        for row in range(4):
            add_shape(slide, I(0.25 + col*0.22), I(6.2 + row*0.22),
                      I(0.08), I(0.08),
                      fill_color=_rgb(0x25,0x63,0xEB), radius=100000)

    # Main title
    add_textbox(slide, "RESUME ENGINE",
                I(0.55), I(1.6), I(9.5), I(1.3),
                font_name="Aptos Display", font_size=Pt(60), bold=True,
                color=WHITE, align=PP_ALIGN.LEFT)

    # Accent underline under title
    add_shape(slide, I(0.55), I(2.82), I(6.5), I(0.06), fill_color=ACCENT_CYAN)

    # Subtitle
    add_textbox(slide, "Git-Driven, Tag-Based Resume Compilation Engine",
                I(0.55), I(2.95), I(9.5), I(0.55),
                font_name="Aptos Display", font_size=Pt(20), bold=False,
                color=ACCENT_CYAN, align=PP_ALIGN.LEFT)

    # Divider
    add_shape(slide, I(0.55), I(3.65), I(4.0), I(0.03), fill_color=_rgb(0x30,0x3B,0x52))

    # Project metadata
    meta_lines = [
        ("Summer Training Project", Pt(14), MUTED),
        ("B.Tech Computer Science — 2025", Pt(13), _rgb(0x64,0x74,0x8B)),
    ]
    for i, (txt, sz, col) in enumerate(meta_lines):
        add_textbox(slide, txt, I(0.55), I(3.8 + i*0.42), I(8), I(0.38),
                    font_name="Aptos", font_size=sz, color=col)

    # Bottom bar
    add_shape(slide, 0, I(7.3), SLIDE_W, I(0.2), fill_color=NAVY_CARD)
    add_textbox(slide, "Presented by: [Author Name]   |   [College Name]   |   [Department]",
                I(0.55), I(7.3), I(10), I(0.2),
                font_size=Pt(9), color=MUTED, align=PP_ALIGN.LEFT)

    add_speaker_notes(slide, """SPEAKER NOTES — Slide 1: Title

WHAT TO SAY:
Good [morning/afternoon]. My name is [Name] and today I'll be presenting Resume Engine — a Git-driven, 
tag-based resume compilation engine I built as part of my Summer Training Project.

The core idea is simple: instead of maintaining multiple hand-edited resume files for different job roles, 
you maintain one YAML file containing your entire career history. You tag every item — every bullet point, 
every project, every skill — with the roles it's relevant to. Then you run one command to generate a 
perfectly tailored, one-page PDF resume for any target role, automatically.

TALKING POINTS:
• This is a real, working software project — not a prototype. It has ~40 unit tests, GitHub Actions CI, 
  and two complete frontends: a CLI and a local web UI.
• The project solves a problem every job seeker faces — the pain of resume customization.
• The title color scheme (#0D1117 background, #2563EB accent) is intentional — it mirrors the actual 
  product's own web UI color scheme.

EXPECTED FACULTY QUESTIONS:
Q: What inspired this project?
A: I noticed I was spending hours copying and pasting resume sections for different job applications 
   and often making them inconsistent. A single source of truth with automatic filtering was the natural solution.
""")
    return slide


def build_slide_02_problem(prs):
    """Slide 2: Problem Statement"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY_BG)
    add_slide_header(slide, "Problem Statement", number=2)

    # 4 problem cards in 2x2 grid
    problems = [
        ("🔁", "Repetitive Manual Work",
         "Customizing a resume per job application means re-editing the same content repeatedly — error-prone and time-consuming."),
        ("🎯", "Role-Specific Emphasis",
         "A backend role needs different highlights than data-science or devops. One-size-fits-all resumes undersell you."),
        ("🔀", "Version Drift",
         "Keeping several hand-edited resume files in sync is easy to get wrong — a fix in one version is forgotten in another."),
        ("📋", "No Single Source of Truth",
         "Career history scattered across multiple resume files with no canonical master record to update once."),
    ]

    cols = [Inches(0.4), Inches(6.85)]
    rows = [Inches(1.1), Inches(4.1)]

    for i, (icon, title, desc) in enumerate(problems):
        col = i % 2
        row = i // 2
        cx, cy = cols[col], rows[row]
        cw, ch = Inches(6.2), Inches(2.7)

        # Card background
        add_card(slide, cx, cy, cw, ch, fill=NAVY_CARD, border=_rgb(0x25,0x3A,0x5E))

        # Left accent bar
        add_shape(slide, cx, cy, Inches(0.07), ch, fill_color=ACCENT_BLUE)

        # Icon
        add_textbox(slide, icon, cx + Inches(0.2), cy + Inches(0.18),
                    Inches(0.7), Inches(0.55), font_size=Pt(28), color=WHITE)

        # Title
        add_textbox(slide, title,
                    cx + Inches(0.9), cy + Inches(0.18), Inches(5.0), Inches(0.45),
                    font_name="Aptos Display", font_size=Pt(16), bold=True, color=WHITE)

        # Description
        add_textbox(slide, desc,
                    cx + Inches(0.2), cy + Inches(0.78), Inches(5.8), Inches(1.75),
                    font_name="Aptos", font_size=Pt(12.5), color=OFF_WHITE, wrap=True)

    # Bottom conclusion line
    add_card(slide, Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.42),
             fill=_rgb(0x1E,0x3A,0x5F), radius=15000)
    add_shape(slide, Inches(0.4), Inches(6.95), Inches(0.06), Inches(0.42),
              fill_color=ACCENT_CYAN)
    add_textbox(slide,
                '  "One master resume should generate unlimited, role-specific resumes automatically."',
                Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.42),
                font_name="Aptos", font_size=Pt(13.5), italic=True, color=ACCENT_CYAN,
                align=PP_ALIGN.LEFT)

    add_speaker_notes(slide, """SPEAKER NOTES — Slide 2: Problem Statement

WHAT TO SAY:
Let me start with the problem. Anyone who has applied to multiple roles knows this pain well.

[Walk through each card]:
1. REPETITIVE WORK: Every time you apply to a different role, you're manually editing the same resume — 
   copying sections, tweaking wording, adding/removing content. It's tedious and error-prone.
2. ROLE-SPECIFIC EMPHASIS: A backend engineer and a data scientist have completely different priorities. 
   A generic resume that tries to cover both ends up being mediocre at communicating either.
3. VERSION DRIFT: The moment you have more than one resume file, they start diverging. A job history 
   update in one file gets missed in another. You end up applying with stale information.
4. NO SOURCE OF TRUTH: There's no canonical record. Your resume lives in fragments across multiple files.

The conclusion: we need a system where career history is written once, and role-specific resumes 
are generated automatically.

EXPECTED FACULTY QUESTIONS:
Q: Isn't this problem solved by LinkedIn or other tools?
A: LinkedIn is a profile, not a formatted PDF resume. Most companies still expect a precisely-formatted 
   one-page PDF, and LinkedIn's export is not suitable. Resume Engine gives you precise PDF control.
""")
    return slide


def build_slide_03_overview(prs):
    """Slide 3: Project Overview"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY_BG)
    add_slide_header(slide, "Project Overview", number=3)

    # One-line goal card
    add_card(slide, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.75),
             fill=_rgb(0x1A,0x2E,0x55), border=ACCENT_BLUE, radius=15000)
    add_textbox(slide,
                "🎯  Maintain ONE YAML source of truth — generate unlimited tailored one-page PDF resumes on demand.",
                Inches(0.6), Inches(1.05), Inches(12.1), Inches(0.75),
                font_name="Aptos", font_size=Pt(14.5), bold=True, color=WHITE,
                align=PP_ALIGN.LEFT)

    # High-level flow diagram
    add_textbox(slide, "High-Level Data Flow",
                Inches(0.4), Inches(2.05), Inches(6), Inches(0.35),
                font_name="Aptos Display", font_size=Pt(15), bold=True,
                color=ACCENT_CYAN)

    # Flow boxes
    flow_items = [
        ("📄", "Tagged\nMaster YAML", NAVY_CARD, ACCENT_BLUE),
        ("⚙️", "Filter Engine\n(tag match +\npriority rank +\nbudget cap)", _rgb(0x1A,0x2E,0x55), ACCENT_BLUE),
        ("📊", "ReportLab\nPDF Renderer", NAVY_CARD, ACCENT_CYAN),
        ("✅", "One-Page\nRole-Specific\nPDF Resume", _rgb(0x14,0x2A,0x1A), GREEN_OK),
    ]

    box_w = Inches(2.6)
    box_h = Inches(2.4)
    gap   = Inches(0.55)
    start_x = Inches(0.4)
    y     = Inches(2.55)

    for i, (icon, label, fill, accent) in enumerate(flow_items):
        x = start_x + i * (box_w + gap)
        add_card(slide, x, y, box_w, box_h, fill=fill, border=accent, radius=25000)
        add_shape(slide, x, y, box_w, Inches(0.06), fill_color=accent)
        # Icon
        add_textbox(slide, icon, x, y + Inches(0.15), box_w, Inches(0.55),
                    font_size=Pt(28), color=WHITE, align=PP_ALIGN.CENTER)
        # Label
        add_textbox(slide, label, x, y + Inches(0.75), box_w, Inches(1.55),
                    font_name="Aptos Display", font_size=Pt(13), bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)

        # Arrow (not after last)
        if i < len(flow_items) - 1:
            ax = x + box_w + Inches(0.1)
            ay = y + box_h/2 - Inches(0.15)
            add_textbox(slide, "  ➜", ax, ay, Inches(0.45), Inches(0.35),
                        font_size=Pt(20), color=ACCENT_CYAN)

    # Two entry points note
    add_textbox(slide, "Two Entry Points:",
                Inches(0.4), Inches(5.2), Inches(3), Inches(0.35),
                font_name="Aptos Display", font_size=Pt(14), bold=True, color=ACCENT_CYAN)

    ep_data = [
        ("🖥️ CLI",    "python compile_resume.py --target backend", ACCENT_BLUE),
        ("🌐 Web UI", "http://127.0.0.1:8000  (FastAPI + Jinja2)",  ACCENT_CYAN),
    ]
    for i, (label, cmd, col) in enumerate(ep_data):
        ex = Inches(0.4) + i * Inches(6.45)
        add_card(slide, ex, Inches(5.65), Inches(6.2), Inches(0.72),
                 fill=NAVY_CARD, radius=15000)
        add_shape(slide, ex, Inches(5.65), Inches(0.07), Inches(0.72), fill_color=col)
        add_textbox(slide, f"{label}  —  {cmd}",
                    ex + Inches(0.2), Inches(5.65), Inches(5.9), Inches(0.72),
                    font_name="Aptos", font_size=Pt(12.5), color=WHITE)

    add_speaker_notes(slide, """SPEAKER NOTES — Slide 3: Project Overview

WHAT TO SAY:
Resume Engine is built around a simple but powerful idea: you write your career history once in a 
tagged YAML file, and the engine does all the heavy lifting of filtering, ranking, and formatting.

Walk the audience through the data flow:
1. Everything starts with the Tagged Master YAML — one file containing ALL your jobs, projects, 
   skills, achievements, each tagged with the roles they're relevant to.
2. The Filter Engine processes this YAML: it keeps only items matching your target role (or tagged 
   "general" for universal items), ranks them by relevance and an optional priority field, and enforces 
   section budgets so the output reliably fits one page.
3. The ReportLab PDF Renderer takes the filtered data and draws directly onto a PDF canvas — no browser, 
   no HTML, no external tools required.
4. The result: a clean, one-page PDF perfectly tailored to your target role.

There are TWO ways to use the engine:
- The CLI for quick command-line compilation
- The Web UI for a browser-based edit-and-generate workflow

Both frontends share the SAME core engine — no duplicated logic. That's a key architectural decision.

EXPECTED FACULTY QUESTIONS:
Q: What's a YAML file?
A: YAML is a human-friendly data format, similar to JSON but much more readable. It's the same format 
   used by tools like Docker Compose and GitHub Actions.
""")
    return slide


def build_slide_04_features(prs):
    """Slide 4: Key Features"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY_BG)
    add_slide_header(slide, "Key Features", accent=ACCENT_CYAN, number=4)

    features = [
        ("🏷️", "Tag-Based Filtering",
         "Items tagged backend, frontend, data-science, devops, general — only the right content gets through."),
        ("⭐", "Priority Ranking",
         "Optional priority: int field gives manual control over item ordering within any section."),
        ("📐", "Section Budgeting",
         "4 bullets/job · 3 projects · 4 achievements · 6 coursework — hard caps ensure reliable one-page output."),
        ("📄", "Multi-Page Overflow",
         "Long profiles spill gracefully to page 2+ instead of being silently cut off."),
        ("⚡", "ReportLab PDF Renderer",
         "Pure Python, no browser, no wkhtmltopdf — fast and dependency-light programmatic rendering."),
        ("🔗", "Git Auto-Commit",
         "Generated PDFs are auto-committed with .gitignore-aware skip logic. Optional --push flag."),
        ("🌐", "Web UI + CLI",
         "Pydantic-validated form editor in the browser; one-command generation from the terminal."),
        ("🧪", "~40 Pytest Tests + CI",
         "Full unit test suite on filter_engine.py; GitHub Actions runs tests + smoke-compiles on every push."),
    ]

    cols = [Inches(0.4), Inches(4.7), Inches(8.95)]
    rows = [Inches(1.1), Inches(2.6), Inches(4.1), Inches(5.6)]

    for i, (icon, title, desc) in enumerate(features):
        if i >= 8:
            break
        row = i // 3
        col = i % 3
        # Last row has only 2 items — center them
        if row == 2 and len(features) % 3 == 2:
            if i == 6:
                cx = Inches(2.5)
            else:
                cx = Inches(6.8)
        else:
            cx = cols[col]
        cy = rows[row]
        fw, fh = Inches(4.1), Inches(1.35)

        add_card(slide, cx, cy, fw, fh, fill=NAVY_CARD, border=_rgb(0x25,0x3A,0x5E), radius=20000)
        add_shape(slide, cx, cy, fw, Inches(0.05), fill_color=ACCENT_BLUE)

        # Icon + title row
        add_textbox(slide, icon,
                    cx + Inches(0.15), cy + Inches(0.1), Inches(0.45), Inches(0.42),
                    font_size=Pt(18), color=WHITE)
        add_textbox(slide, title,
                    cx + Inches(0.6), cy + Inches(0.1), Inches(3.35), Inches(0.42),
                    font_name="Aptos Display", font_size=Pt(13), bold=True, color=WHITE)

        # Description
        add_textbox(slide, desc,
                    cx + Inches(0.15), cy + Inches(0.6), Inches(3.8), Inches(0.72),
                    font_name="Aptos", font_size=Pt(10.5), color=OFF_WHITE, wrap=True)

    add_speaker_notes(slide, """SPEAKER NOTES — Slide 4: Key Features

WHAT TO SAY:
These are the eight core features that are FULLY IMPLEMENTED and working. I'll call out a few key ones:

TAG-BASED FILTERING: This is the heart of the system. Every item in your YAML carries a "tags" list. 
The engine keeps items tagged for your target role, plus anything tagged "general" (universal items).

PRIORITY RANKING: This is a thoughtful UX addition — you can add priority: 2 to a bullet point to force 
it to the top of a section, regardless of relevance score. Completely optional.

SECTION BUDGETING: Hard numerical caps — 4 bullets per job, max 3 projects, etc. This is what guarantees 
the one-page output. It's enforced in apply_budget() in filter_engine.py.

GIT AUTO-COMMIT: Every time you generate a PDF, the engine can auto-commit it to your git history. 
This gives you a version trail of your resume changes over time. There's a safety: if the PDF is 
gitignored (e.g., in the output/ directory), the commit is silently skipped.

~40 TESTS: These tests cover the filtering, ranking, and budget logic comprehensively with edge cases 
like empty lists, missing tags, and unknown targets.

IMPORTANT HONESTY NOTE: The web UI has a "Theme" dropdown that says "Coming Soon" — the renderer 
currently has one fixed design. Multiple themes are on the roadmap. I'll cover this on the Future Scope slide.

EXPECTED FACULTY QUESTIONS:
Q: What does "general" tag mean?
A: Items tagged "general" appear in ALL target resumes — things like your contact info, or a skill 
   that's relevant to every role you'd apply to.
""")
    return slide


def build_slide_05_techstack(prs):
    """Slide 5: Technology Stack"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY_BG)
    add_slide_header(slide, "Technology Stack", number=5)

    stack = [
        ("Language",        "Python 3.12",                       "🐍", ACCENT_BLUE),
        ("Data Format",     "YAML  (PyYAML)",                    "📦", ACCENT_CYAN),
        ("PDF Generation",  "ReportLab (canvas API)",            "📄", _rgb(0xF5,0x9E,0x0B)),
        ("Web Framework",   "FastAPI + Uvicorn (ASGI)",          "⚡", _rgb(0x10,0xB9,0x81)),
        ("Templating",      "Jinja2",                            "🔧", _rgb(0xA7,0x8B,0xFA)),
        ("Data Validation", "Pydantic  (web layer only)",        "✅", _rgb(0xFB,0x92,0x3C)),
        ("Frontend",        "HTML5 · CSS3 · Vanilla JavaScript", "🌐", ACCENT_CYAN),
        ("Git Automation",  "GitPython",                         "🔗", _rgb(0xF4,0x73,0x73)),
        ("CLI UX",          "argparse · colorama",               "🖥️", ACCENT_BLUE),
        ("Testing",         "pytest  (~40 unit tests)",          "🧪", _rgb(0x22,0xC5,0x5E)),
        ("CI/CD",           "GitHub Actions",                    "🚀", _rgb(0xFB,0xD5,0x68)),
    ]

    cols = 3
    rows_per_col = 4
    sw = Inches(4.1)
    sh = Inches(0.8)
    sg_x = Inches(0.3)
    sg_y = Inches(0.22)
    start_x = Inches(0.35)
    start_y = Inches(1.15)

    for i, (cat, tech, icon, col) in enumerate(stack):
        ci = i % cols
        ri = i // cols
        sx = start_x + ci * (sw + sg_x)
        sy = start_y + ri * (sh + sg_y)

        add_card(slide, sx, sy, sw, sh, fill=NAVY_CARD, radius=15000)
        # Left colored accent
        add_shape(slide, sx, sy, Inches(0.06), sh, fill_color=col)

        # Icon
        add_textbox(slide, icon, sx + Inches(0.15), sy, Inches(0.4), sh,
                    font_size=Pt(18), color=WHITE, align=PP_ALIGN.CENTER)

        # Category label
        add_textbox(slide, cat.upper(),
                    sx + Inches(0.58), sy + Inches(0.06), Inches(3.3), Inches(0.28),
                    font_name="Aptos", font_size=Pt(8.5), bold=True, color=MUTED)

        # Tech name
        add_textbox(slide, tech,
                    sx + Inches(0.58), sy + Inches(0.35), Inches(3.3), Inches(0.4),
                    font_name="Aptos Display", font_size=Pt(13), bold=True, color=WHITE)

    add_speaker_notes(slide, """SPEAKER NOTES — Slide 5: Technology Stack

WHAT TO SAY:
The tech stack was chosen for pragmatism and minimal external dependencies.

Let me highlight the key choices:

PYTHON 3.12: The CI config explicitly pins Python 3.12, so this is precise, not approximate.

REPORTLAB: This is the most unusual choice — most people would reach for an HTML renderer. 
I chose ReportLab because it has zero external dependencies (no browser, no wkhtmltopdf). 
The tradeoff is that layout control requires explicit coordinate-based drawing, which was a 
genuine engineering challenge — but it makes the system much easier to deploy.

FASTAPI: A modern, async Python web framework with automatic API documentation. Much faster and 
cleaner than Flask for this use case, especially with Pydantic integration.

PYDANTIC: Critically, Pydantic validation only happens on the WEB EDITOR's save path. The CLI 
compile path works with plain Python dicts throughout — this is an important architectural detail.

NO FRONTEND FRAMEWORK: The web UI is plain HTML/CSS/JavaScript — no React, no Vue. This was 
a deliberate choice to keep the project focused on the backend engineering problem.

EXPECTED FACULTY QUESTIONS:
Q: Why not use a database?
A: The project's data model is a single user's career history — a YAML file is perfectly appropriate 
   for this. A database would add complexity with no benefit. It IS a future roadmap item if 
   multi-user support were ever added.

Q: Why Python and not something faster?
A: PDF generation is I/O bound, not compute bound. Python's ecosystem (ReportLab, FastAPI, GitPython) 
   is exactly right for this type of document-generation tool.
""")
    return slide


def build_slide_06_architecture(prs):
    """Slide 6: High-Level Architecture"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY_BG)
    add_slide_header(slide, "High-Level Architecture", number=6)

    # ── Entry points column ──────────────────────────────────────────────────
    add_textbox(slide, "ENTRY POINTS",
                Inches(0.4), Inches(1.1), Inches(3.5), Inches(0.32),
                font_size=Pt(10), bold=True, color=MUTED)

    cli_y = Inches(1.5)
    web_y = Inches(3.2)

    # CLI box
    add_card(slide, Inches(0.4), cli_y, Inches(3.4), Inches(1.45),
             fill=_rgb(0x1A,0x2E,0x55), border=ACCENT_BLUE, radius=20000)
    add_shape(slide, Inches(0.4), cli_y, Inches(3.4), Inches(0.07), fill_color=ACCENT_BLUE)
    add_textbox(slide, "🖥️  CLI",
                Inches(0.5), cli_y + Inches(0.1), Inches(3.2), Inches(0.45),
                font_name="Aptos Display", font_size=Pt(16), bold=True, color=ACCENT_BLUE)
    add_textbox(slide, "compile_resume.py",
                Inches(0.5), cli_y + Inches(0.58), Inches(3.2), Inches(0.35),
                font_name="Aptos", font_size=Pt(11), color=OFF_WHITE)
    add_textbox(slide, "python compile_resume.py --target backend",
                Inches(0.5), cli_y + Inches(0.9), Inches(3.2), Inches(0.4),
                font_name="Aptos", font_size=Pt(9), italic=True, color=MUTED)

    # Web UI box
    add_card(slide, Inches(0.4), web_y, Inches(3.4), Inches(1.45),
             fill=_rgb(0x14,0x2A,0x30), border=ACCENT_CYAN, radius=20000)
    add_shape(slide, Inches(0.4), web_y, Inches(3.4), Inches(0.07), fill_color=ACCENT_CYAN)
    add_textbox(slide, "🌐  Web UI",
                Inches(0.5), web_y + Inches(0.1), Inches(3.2), Inches(0.45),
                font_name="Aptos Display", font_size=Pt(16), bold=True, color=ACCENT_CYAN)
    add_textbox(slide, "web/app.py  (FastAPI)",
                Inches(0.5), web_y + Inches(0.58), Inches(3.2), Inches(0.35),
                font_name="Aptos", font_size=Pt(11), color=OFF_WHITE)
    add_textbox(slide, "Browser → GET / · POST /save · POST /generate",
                Inches(0.5), web_y + Inches(0.9), Inches(3.2), Inches(0.4),
                font_name="Aptos", font_size=Pt(9), italic=True, color=MUTED)

    # User box
    add_card(slide, Inches(0.4), Inches(5.0), Inches(3.4), Inches(0.7),
             fill=NAVY_CARD, radius=20000)
    add_textbox(slide, "👤  User  (CLI or Browser)",
                Inches(0.5), Inches(5.0), Inches(3.2), Inches(0.7),
                font_name="Aptos", font_size=Pt(13), bold=True, color=OFF_WHITE,
                align=PP_ALIGN.CENTER)

    # Arrows from User to entry points
    for y in [Inches(2.28), Inches(4.0)]:
        add_textbox(slide, "▼", Inches(1.8), y, Inches(0.35), Inches(0.32),
                    font_size=Pt(14), color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    # ── Shared Core column ───────────────────────────────────────────────────
    add_textbox(slide, "SHARED CORE  (same code for both frontends)",
                Inches(4.3), Inches(1.1), Inches(5.7), Inches(0.32),
                font_size=Pt(10), bold=True, color=MUTED)

    core_items = [
        ("⚙️  filter_engine.py",
         "YAML loading → tag filter → rank → budget cap",
         Inches(4.3), Inches(1.5), ACCENT_BLUE),
        ("📄  pdf_renderer.py",
         "ReportLab canvas renderer · multi-page overflow",
         Inches(4.3), Inches(2.9), ACCENT_CYAN),
        ("🔗  git_utils.py",
         "Auto-commit PDF · optional push · .gitignore-aware",
         Inches(4.3), Inches(4.3), _rgb(0x22,0xC5,0x5E)),
    ]
    for (title, sub, cx, cy, col) in core_items:
        add_card(slide, cx, cy, Inches(5.5), Inches(1.1),
                 fill=NAVY_CARD, border=col, radius=20000)
        add_shape(slide, cx, cy, Inches(5.5), Inches(0.06), fill_color=col)
        add_textbox(slide, title, cx + Inches(0.15), cy + Inches(0.08),
                    Inches(5.2), Inches(0.42),
                    font_name="Aptos Display", font_size=Pt(15), bold=True, color=WHITE)
        add_textbox(slide, sub, cx + Inches(0.15), cy + Inches(0.55),
                    Inches(5.2), Inches(0.45),
                    font_name="Aptos", font_size=Pt(11.5), color=OFF_WHITE)

    # Connecting arrows from entry points to core
    for ay in [Inches(2.22), Inches(3.92)]:
        add_textbox(slide, "→", Inches(3.82), ay, Inches(0.45), Inches(0.35),
                    font_size=Pt(18), color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

    # Down arrow in core (filter → renderer → git)
    for ay in [Inches(2.62), Inches(4.02)]:
        add_textbox(slide, "↓", Inches(7.1), ay, Inches(0.35), Inches(0.32),
                    font_size=Pt(16), color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

    # ── Output column ────────────────────────────────────────────────────────
    add_textbox(slide, "OUTPUT",
                Inches(10.3), Inches(1.1), Inches(2.7), Inches(0.32),
                font_size=Pt(10), bold=True, color=MUTED)

    add_card(slide, Inches(10.3), Inches(2.35), Inches(2.7), Inches(1.2),
             fill=_rgb(0x14,0x2A,0x1A), border=GREEN_OK, radius=20000)
    add_textbox(slide, "✅  PDF Output",
                Inches(10.35), Inches(2.45), Inches(2.55), Inches(0.4),
                font_name="Aptos Display", font_size=Pt(14), bold=True, color=GREEN_OK,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, "output/resume_{target}.pdf",
                Inches(10.35), Inches(2.9), Inches(2.55), Inches(0.35),
                font_name="Aptos", font_size=Pt(9.5), color=OFF_WHITE,
                align=PP_ALIGN.CENTER, italic=True)

    add_textbox(slide, "→", Inches(9.82), Inches(2.75), Inches(0.45), Inches(0.35),
                font_size=Pt(18), color=GREEN_OK, align=PP_ALIGN.CENTER)

    # Shared core callout box at bottom
    add_card(slide, Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.72),
             fill=_rgb(0x1A,0x2E,0x55), border=ACCENT_BLUE, radius=15000)
    add_textbox(slide,
                "★  Both frontends share the exact same filter_engine.py and pdf_renderer.py — zero logic duplication.",
                Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.72),
                font_name="Aptos", font_size=Pt(13), bold=True, color=ACCENT_CYAN,
                align=PP_ALIGN.LEFT)

    add_speaker_notes(slide, """SPEAKER NOTES — Slide 6: High-Level Architecture

WHAT TO SAY:
This is the most important architectural slide. The key insight is the two-frontend, one-core design.

On the LEFT are the two entry points: the CLI (compile_resume.py) and the Web UI (web/app.py).
In the MIDDLE is the shared core: filter_engine, pdf_renderer, and git_utils.
On the RIGHT is the output: the generated PDF.

The CRITICAL POINT to emphasize:
Both frontends — the CLI and the web UI — call the EXACT SAME filter_engine.py and pdf_renderer.py. 
There is zero duplication of the core logic. This is a genuine separation-of-concerns decision.

What this means in practice:
• A bug fix in filter_engine.py benefits both the CLI and the Web UI simultaneously.
• Adding a new feature to the renderer automatically works in both frontends.
• The core is independently testable — and it IS independently tested with ~40 unit tests.

The git_utils module is optional — it only runs when requested via the --push flag or web UI checkbox.

EXPECTED FACULTY QUESTIONS:
Q: Why separate the CLI and Web UI instead of just having one interface?
A: Different use cases. The CLI is faster for developers who know what target they want. 
   The Web UI is better for editing resume data and for users who prefer a visual interface. 
   The shared core means building both doesn't cost much — the hard work is in the core.
""")
    return slide


def build_slide_07_modules(prs):
    """Slide 7: Module Architecture"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY_BG)
    add_slide_header(slide, "Module Architecture", number=7)

    # Layer 1: Frontends
    add_textbox(slide, "PRESENTATION LAYER  (Entry Points)",
                Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.3),
                font_size=Pt(9), bold=True, color=MUTED)

    modules_l1 = [
        ("compile_resume.py", "CLI Entry Point\nargparse · colorama\n--target, --output, --push, --preview\n--list-targets, --targets-report",
         Inches(0.4), ACCENT_BLUE),
        ("web/app.py", "Web Entry Point\nFastAPI routes\nGET /  ·  POST /save\nGET+POST /generate  ·  POST /git/push",
         Inches(5.0), ACCENT_CYAN),
    ]
    for (mod, desc, x, col) in modules_l1:
        add_card(slide, x, Inches(1.42), Inches(4.45), Inches(1.4),
                 fill=NAVY_CARD, border=col, radius=18000)
        add_shape(slide, x, Inches(1.42), Inches(4.45), Inches(0.06), fill_color=col)
        add_textbox(slide, mod, x + Inches(0.15), Inches(1.5),
                    Inches(4.2), Inches(0.38),
                    font_name="Aptos Display", font_size=Pt(14), bold=True, color=col)
        add_textbox(slide, desc, x + Inches(0.15), Inches(1.9),
                    Inches(4.2), Inches(0.85),
                    font_name="Aptos", font_size=Pt(9.5), color=OFF_WHITE, wrap=True)

    # Web-only side branch
    add_card(slide, Inches(9.7), Inches(1.42), Inches(3.3), Inches(2.7),
             fill=_rgb(0x1A,0x20,0x30), border=_rgb(0x38,0x45,0x65), radius=18000)
    add_textbox(slide, "WEB-ONLY LAYER",
                Inches(9.8), Inches(1.48), Inches(3.1), Inches(0.28),
                font_size=Pt(9), bold=True, color=MUTED)
    web_mods = [
        ("web/models.py", "Pydantic schemas\nResumeData + nested models\nValidation on save path only"),
        ("web/data_store.py", "YAML file I/O\nSingle choke-point for\nall web data reads/writes"),
        ("web/templates/ + static/", "Jinja2 HTML · CSS dark theme\nVanilla JS — no framework"),
    ]
    for j, (mod, desc) in enumerate(web_mods):
        add_textbox(slide, f"▸ {mod}",
                    Inches(9.85), Inches(1.8) + j*Inches(0.78), Inches(3.0), Inches(0.3),
                    font_name="Aptos", font_size=Pt(10.5), bold=True, color=ACCENT_CYAN)
        add_textbox(slide, desc,
                    Inches(9.85), Inches(2.06) + j*Inches(0.78), Inches(3.0), Inches(0.5),
                    font_name="Aptos", font_size=Pt(9), color=OFF_WHITE, wrap=True)

    # Arrows from L1 → L2
    for ax in [Inches(2.6), Inches(7.2)]:
        add_textbox(slide, "↓", ax, Inches(2.88), Inches(0.3), Inches(0.3),
                    font_size=Pt(14), color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

    # Layer 2: Core engine
    add_textbox(slide, "CORE ENGINE  (Shared — no duplication)",
                Inches(0.4), Inches(3.22), Inches(9.0), Inches(0.28),
                font_size=Pt(9), bold=True, color=MUTED)

    core_mods = [
        ("filter_engine.py", "load_yaml() → filter_list() → rank_by_relevance()\n→ apply_budget() → build_resume_data()\nPure functions — no side effects — fully unit tested",
         Inches(0.4), ACCENT_BLUE, _rgb(0x1A,0x2E,0x55)),
        ("pdf_renderer.py", "Renderer class — ReportLab canvas\n_check_or_new_page() overflow gate\nProgrammatic drawing — no HTML/browser",
         Inches(4.8), ACCENT_CYAN, _rgb(0x14,0x24,0x30)),
    ]
    for (mod, desc, x, col, fill) in core_mods:
        add_card(slide, x, Inches(3.52), Inches(4.2), Inches(1.5),
                 fill=fill, border=col, radius=18000)
        add_shape(slide, x, Inches(3.52), Inches(4.2), Inches(0.06), fill_color=col)
        add_textbox(slide, mod, x + Inches(0.15), Inches(3.6),
                    Inches(4.0), Inches(0.38),
                    font_name="Aptos Display", font_size=Pt(14), bold=True, color=col)
        add_textbox(slide, desc, x + Inches(0.15), Inches(4.0),
                    Inches(4.0), Inches(0.95),
                    font_name="Aptos", font_size=Pt(9.5), color=OFF_WHITE, wrap=True)

    # Arrow between core modules
    add_textbox(slide, "→", Inches(4.42), Inches(4.1), Inches(0.35), Inches(0.3),
                font_size=Pt(14), color=ACCENT_CYAN)

    # Arrow to git_utils
    add_textbox(slide, "↓", Inches(6.8), Inches(5.08), Inches(0.3), Inches(0.28),
                font_size=Pt(14), color=GREEN_OK, align=PP_ALIGN.CENTER)

    # git_utils
    add_textbox(slide, "OPTIONAL LAYER",
                Inches(0.4), Inches(5.4), Inches(9.0), Inches(0.28),
                font_size=Pt(9), bold=True, color=MUTED)
    add_card(slide, Inches(4.8), Inches(5.42), Inches(4.2), Inches(1.05),
             fill=_rgb(0x14,0x2A,0x1A), border=GREEN_OK, radius=18000)
    add_shape(slide, Inches(4.8), Inches(5.42), Inches(4.2), Inches(0.06), fill_color=GREEN_OK)
    add_textbox(slide, "git_utils.py",
                Inches(4.95), Inches(5.5), Inches(3.9), Inches(0.38),
                font_name="Aptos Display", font_size=Pt(14), bold=True, color=GREEN_OK)
    add_textbox(slide, "git_commit() · git_push() · .gitignore-aware · graceful no-op outside git repo",
                Inches(4.95), Inches(5.9), Inches(3.9), Inches(0.5),
                font_name="Aptos", font_size=Pt(9.5), color=OFF_WHITE, wrap=True)

    # tests
    add_card(slide, Inches(0.4), Inches(5.42), Inches(4.1), Inches(1.05),
             fill=_rgb(0x14,0x20,0x14), border=_rgb(0x22,0xC5,0x5E), radius=18000)
    add_textbox(slide, "tests/test_filter_engine.py",
                Inches(0.55), Inches(5.5), Inches(3.8), Inches(0.38),
                font_name="Aptos Display", font_size=Pt(12), bold=True, color=GREEN_OK)
    add_textbox(slide, "~40 pytest tests · covers filter_engine only\nConftest.py with shared fixtures",
                Inches(0.55), Inches(5.9), Inches(3.8), Inches(0.5),
                font_name="Aptos", font_size=Pt(9.5), color=OFF_WHITE, wrap=True)

    # Bottom note
    add_card(slide, Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.65),
             fill=_rgb(0x1A,0x26,0x1A), radius=12000)
    add_textbox(slide,
                "⚠  Pydantic validation (web/models.py) applies only on the web /save path — the CLI compile path uses plain Python dicts throughout.",
                Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.65),
                font_name="Aptos", font_size=Pt(11), italic=True, color=AMBER,
                align=PP_ALIGN.LEFT)

    add_speaker_notes(slide, """SPEAKER NOTES — Slide 7: Module Architecture

WHAT TO SAY:
Let me walk through the module layers from top to bottom.

PRESENTATION LAYER: Two entry points. compile_resume.py handles the CLI with argparse, colorama 
for colored output, and implements flags like --target, --output, --push, --preview, 
--list-targets, and --targets-report. web/app.py is the FastAPI application with five routes.

WEB-ONLY LAYER (right side): Three modules that are ONLY used in the web path:
- web/models.py: Pydantic schemas that validate data when you save through the web editor.
- web/data_store.py: A single choke-point for all YAML file I/O — this is a clean design 
  that prevents scattered file access across the web layer.
- Templates and static files: Jinja2 HTML with a dark-theme CSS and vanilla JavaScript.

CORE ENGINE (bottom): The two workhorse modules that both frontends share:
- filter_engine.py: Pure functions — load, filter, rank, budget, assemble. No side effects.
- pdf_renderer.py: A Renderer class that draws directly onto a ReportLab canvas.

IMPORTANT CLARIFICATION: Pydantic validation happens ONLY in the web save path. The CLI compile 
path never goes through Pydantic — data stays as plain Python dicts, accessed with .get() and defaults.

EXPECTED FACULTY QUESTIONS:
Q: Why is web/data_store.py a separate module?
A: Single responsibility principle. Centralizing all YAML I/O in one module means if we ever need 
   to change the storage backend (e.g., to a database), we only change one file.
""")
    return slide


def build_slide_08_dataflow(prs):
    """Slide 8: Data Flow Diagram"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY_BG)
    add_slide_header(slide, "Data Flow — Compile Path", number=8)

    # Two tracks: CLI and Web, but we show the shared compile path in detail
    steps = [
        ("📂", "data/resume.yaml\n(or .example.yaml fallback)",
         "YAML source — gitignored real data,\nor fictional example for demo/CI",
         NAVY_CARD, ACCENT_BLUE),
        ("⚙️", "filter_engine.\nload_yaml()",
         "yaml.safe_load() → plain Python dict\nNo Pydantic, no schema enforcement",
         _rgb(0x1A,0x2E,0x55), ACCENT_BLUE),
        ("🔽", "Per-section\nFilter → Rank → Budget",
         "filter_list()  →  rank_by_relevance()\n→  apply_budget()  per section",
         _rgb(0x14,0x24,0x34), ACCENT_CYAN),
        ("🔨", "build_resume\n_data()",
         "Assembles filtered dict:\npersonal · education · experience\nprojects · skills · achievements",
         NAVY_CARD, ACCENT_CYAN),
        ("📄", "pdf_renderer.\nrender_pdf()",
         "Renderer draws onto ReportLab Canvas\n_check_or_new_page() handles overflow",
         _rgb(0x14,0x20,0x2E), GREEN_OK),
        ("💾", "output/resume\n_{target}.pdf",
         "Written to disk\nMulti-page if content overflows",
         _rgb(0x10,0x28,0x18), GREEN_OK),
        ("🔗", "git_utils.\ngit_commit()\n(optional)",
         "Auto-commits PDF if not gitignored\ngit_push() only if --push flag",
         _rgb(0x1A,0x2A,0x1A), _rgb(0x22,0xC5,0x5E)),
    ]

    bw = Inches(1.62)
    bh = Inches(2.0)
    gap = Inches(0.25)
    start_x = Inches(0.35)
    y = Inches(1.15)

    for i, (icon, title, desc, fill, col) in enumerate(steps):
        x = start_x + i * (bw + gap)
        add_card(slide, x, y, bw, bh, fill=fill, border=col, radius=18000)
        add_shape(slide, x, y, bw, Inches(0.055), fill_color=col)

        # Icon
        add_textbox(slide, icon, x, y + Inches(0.08), bw, Inches(0.4),
                    font_size=Pt(18), color=WHITE, align=PP_ALIGN.CENTER)

        # Title
        add_textbox(slide, title, x + Inches(0.05), y + Inches(0.5), bw - Inches(0.1), Inches(0.72),
                    font_name="Aptos Display", font_size=Pt(10.5), bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)

        # Desc
        add_textbox(slide, desc, x + Inches(0.06), y + Inches(1.25), bw - Inches(0.12), Inches(0.68),
                    font_name="Aptos", font_size=Pt(8.5), color=OFF_WHITE,
                    align=PP_ALIGN.LEFT, wrap=True)

        # Arrow (not after last)
        if i < len(steps) - 1:
            ax = x + bw + Inches(0.04)
            add_textbox(slide, "➜", ax, y + bh/2 - Inches(0.15), Inches(0.22), Inches(0.3),
                        font_size=Pt(11), color=ACCENT_CYAN)

    # Budget constants callout
    add_card(slide, Inches(0.35), Inches(3.35), Inches(12.6), Inches(1.35),
             fill=_rgb(0x1A,0x20,0x30), border=_rgb(0x30,0x45,0x6A), radius=15000)
    add_textbox(slide, "Section Budget Constants (enforced by apply_budget() in filter_engine.py):",
                Inches(0.55), Inches(3.42), Inches(12.0), Inches(0.32),
                font_name="Aptos", font_size=Pt(11), bold=True, color=ACCENT_CYAN)

    budgets = [
        ("MAX_BULLETS_PER_JOB", "4"),
        ("MAX_PROJECTS", "3"),
        ("MAX_ACHIEVEMENTS", "4"),
        ("MAX_COURSEWORK", "6"),
    ]
    for i, (const, val) in enumerate(budgets):
        bx = Inches(0.55) + i * Inches(3.15)
        add_textbox(slide, f"{const}  =  {val}",
                    bx, Inches(3.82), Inches(3.0), Inches(0.3),
                    font_name="Aptos", font_size=Pt(11), color=OFF_WHITE)

    # Important note: no Pydantic on this path
    add_card(slide, Inches(0.35), Inches(4.85), Inches(12.6), Inches(0.55),
             fill=_rgb(0x28,0x20,0x10), border=AMBER, radius=12000)
    add_textbox(slide,
                "⚠  No Pydantic validation on the CLI compile path — data stays as plain Python dicts, accessed with .get() and defaults throughout.",
                Inches(0.55), Inches(4.85), Inches(12.2), Inches(0.55),
                font_name="Aptos", font_size=Pt(11), italic=True, color=AMBER)

    # Web path note
    add_card(slide, Inches(0.35), Inches(5.6), Inches(12.6), Inches(0.85),
             fill=_rgb(0x14,0x22,0x30), border=ACCENT_CYAN, radius=12000)
    add_textbox(slide, "Web Path Addition:",
                Inches(0.55), Inches(5.65), Inches(3.0), Inches(0.32),
                font_name="Aptos", font_size=Pt(11), bold=True, color=ACCENT_CYAN)
    add_textbox(slide,
                "Browser  →  GET /  →  Edit form (pre-filled from YAML)  →  POST /save  (Pydantic validates → writes YAML)  →  POST /generate  →  same filter_engine + pdf_renderer  →  PDF download",
                Inches(0.55), Inches(5.97), Inches(12.2), Inches(0.42),
                font_name="Aptos", font_size=Pt(10.5), color=OFF_WHITE, wrap=True)

    add_speaker_notes(slide, """SPEAKER NOTES — Slide 8: Data Flow

WHAT TO SAY:
This slide shows exactly what happens when you run compile_resume.py or click Generate in the web UI.

Step 1: YAML Source. The engine first resolves which YAML to use — data/resume.yaml if it exists, 
falling back to data/resume.example.yaml. This fallback is what makes the CI pipeline work safely 
without anyone's personal data being committed.

Step 2: load_yaml(). yaml.safe_load() turns the YAML into a plain Python dict. NO Pydantic here 
— this is important, because Pydantic validation only happens in the web editor's save path.

Step 3: Per-section processing. For EACH section (education, experience, projects, skills, achievements):
  - filter_list(): keeps items tagged for the target or "general"  
  - rank_by_relevance(): sorts by priority field first, then tag count
  - apply_budget(): trims to the hard cap

Step 4: build_resume_data() assembles all filtered sections into one dict.

Step 5: render_pdf() draws everything onto a ReportLab canvas. The _check_or_new_page() method 
handles overflow — if content would go past the page bottom, it starts a new page automatically.

Step 6: PDF written to disk.

Step 7 (optional): git_commit() — if inside a git repo and the file isn't gitignored.

IMPORTANT: Show the budget constants card. These specific numbers (4, 3, 4, 6) are what make 
"fits on one page" a reliable guarantee.

EXPECTED FACULTY QUESTIONS:
Q: What if the content still doesn't fit one page even with budgets?
A: The renderer's _check_or_new_page() starts a new page rather than truncating content silently. 
   So long profiles can produce 2-page resumes — the budget constants are tuned to make one-page 
   the common case.
""")
    return slide


def build_slide_09_workflow(prs):
    """Slide 9: Project Workflow"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY_BG)
    add_slide_header(slide, "Project Workflow", number=9)

    # Two horizontal workflow tracks
    # ─── CLI Track ──────────────────────────────────────────────────────────
    add_textbox(slide, "🖥️  CLI Workflow",
                Inches(0.4), Inches(1.1), Inches(6.3), Inches(0.38),
                font_name="Aptos Display", font_size=Pt(16), bold=True, color=ACCENT_BLUE)
    add_shape(slide, Inches(0.4), Inches(1.5), Inches(5.8), Inches(0.04),
              fill_color=ACCENT_BLUE)

    cli_steps = [
        ("Run Command", "python compile_resume.py\n--target backend"),
        ("Load YAML", "Resolve data/resume.yaml\nor .example fallback"),
        ("Filter/Rank/Budget", "filter_engine per section\n→ build_resume_data()"),
        ("Render PDF", "pdf_renderer\n→ output/resume_backend.pdf"),
        ("Git Commit\n(optional)", "git_utils.git_commit()\ngit_utils.git_push()"),
    ]

    cli_y = Inches(1.65)
    sw = Inches(2.2)
    sh = Inches(1.4)
    gap = Inches(0.2)
    for i, (title, desc) in enumerate(cli_steps):
        x = Inches(0.4) + i * (sw + gap)
        add_card(slide, x, cli_y, sw, sh, fill=NAVY_CARD, border=ACCENT_BLUE, radius=18000)
        add_shape(slide, x, cli_y, sw, Inches(0.05), fill_color=ACCENT_BLUE)
        add_textbox(slide, f"{i+1}. {title}",
                    x + Inches(0.12), cli_y + Inches(0.08), sw - Inches(0.15), Inches(0.45),
                    font_name="Aptos Display", font_size=Pt(11.5), bold=True, color=WHITE)
        add_textbox(slide, desc,
                    x + Inches(0.12), cli_y + Inches(0.55), sw - Inches(0.15), Inches(0.78),
                    font_name="Aptos", font_size=Pt(9.5), color=OFF_WHITE, wrap=True)
        if i < len(cli_steps) - 1:
            ax = x + sw + Inches(0.02)
            add_textbox(slide, "→", ax, cli_y + sh/2 - Inches(0.15), Inches(0.18), Inches(0.3),
                        font_size=Pt(11), color=ACCENT_BLUE)

    # ─── Web Track ──────────────────────────────────────────────────────────
    add_textbox(slide, "🌐  Web UI Workflow",
                Inches(0.4), Inches(3.3), Inches(6.3), Inches(0.38),
                font_name="Aptos Display", font_size=Pt(16), bold=True, color=ACCENT_CYAN)
    add_shape(slide, Inches(0.4), Inches(3.7), Inches(5.8), Inches(0.04),
              fill_color=ACCENT_CYAN)

    web_steps = [
        ("Open Editor", "GET /\nPre-filled from YAML\nPydantic-normalized"),
        ("Edit & Save", "POST /save\nPydantic validates\nWrites data/resume.yaml"),
        ("Choose Target", "GET /generate\nPick tag (backend, etc.)\nOptional: commit to git ☑"),
        ("Generate PDF", "POST /generate\nSame filter_engine\n+ pdf_renderer core"),
        ("Download\n+ Optional Push", "PDF streamed back\nPOST /git/push\n(confirm() dialog)"),
    ]

    web_y = Inches(3.85)
    for i, (title, desc) in enumerate(web_steps):
        x = Inches(0.4) + i * (sw + gap)
        add_card(slide, x, web_y, sw, sh, fill=NAVY_CARD, border=ACCENT_CYAN, radius=18000)
        add_shape(slide, x, web_y, sw, Inches(0.05), fill_color=ACCENT_CYAN)
        add_textbox(slide, f"{i+1}. {title}",
                    x + Inches(0.12), web_y + Inches(0.08), sw - Inches(0.15), Inches(0.45),
                    font_name="Aptos Display", font_size=Pt(11.5), bold=True, color=WHITE)
        add_textbox(slide, desc,
                    x + Inches(0.12), web_y + Inches(0.55), sw - Inches(0.15), Inches(0.78),
                    font_name="Aptos", font_size=Pt(9.5), color=OFF_WHITE, wrap=True)
        if i < len(web_steps) - 1:
            ax = x + sw + Inches(0.02)
            add_textbox(slide, "→", ax, web_y + sh/2 - Inches(0.15), Inches(0.18), Inches(0.3),
                        font_size=Pt(11), color=ACCENT_CYAN)

    # Convergence note
    add_card(slide, Inches(0.4), Inches(5.55), Inches(12.5), Inches(0.55),
             fill=_rgb(0x1A,0x2E,0x55), border=ACCENT_BLUE, radius=12000)
    add_textbox(slide,
                "⚡  Both tracks converge at filter_engine.py + pdf_renderer.py — the same shared core executes regardless of entry point.",
                Inches(0.6), Inches(5.55), Inches(12.1), Inches(0.55),
                font_name="Aptos", font_size=Pt(12), bold=True, color=ACCENT_CYAN)

    # CLI commands reference
    add_card(slide, Inches(0.4), Inches(6.25), Inches(12.5), Inches(1.05),
             fill=NAVY_CARD, radius=12000)
    add_textbox(slide, "Quick CLI Reference:",
                Inches(0.6), Inches(6.3), Inches(3), Inches(0.3),
                font_size=Pt(10), bold=True, color=MUTED)
    cli_ref = "python compile_resume.py --target backend      │      --push (auto-commit+push)      │      --preview (open PDF)      │      --list-targets      │      --targets-report"
    add_textbox(slide, cli_ref,
                Inches(0.6), Inches(6.62), Inches(12.1), Inches(0.55),
                font_name="Aptos", font_size=Pt(10.5), italic=True, color=OFF_WHITE)

    add_speaker_notes(slide, """SPEAKER NOTES — Slide 9: Project Workflow

WHAT TO SAY:
This slide shows the two workflows side-by-side. Let me walk through each.

CLI WORKFLOW (top track):
1. You run python compile_resume.py --target backend
2. The engine resolves the YAML file — your personal data if it exists, demo data otherwise
3. filter_engine processes each section: filter, rank by relevance/priority, apply budget caps
4. pdf_renderer draws the result onto a ReportLab canvas → writes the PDF
5. Optionally, git_utils auto-commits the PDF (skips safely if gitignored or outside a git repo)

WEB WORKFLOW (bottom track):
1. Open browser → GET / shows the Edit Resume page, pre-filled from your YAML
2. You edit fields and click Save → POST /save runs Pydantic validation and writes the YAML
3. Navigate to GET /generate → pick your target role and optionally check "commit to git"
4. Click Generate → POST /generate calls the SAME filter_engine + pdf_renderer as the CLI
5. PDF streams back as a download. Optional separate step: POST /git/push with a browser confirm() dialog

KEY POINT: Steps 4 on the web track and steps 3-4 on the CLI track are IDENTICAL in terms of the 
code that runs. Different input path, same core.

EXPECTED FACULTY QUESTIONS:
Q: What happens if the web server crashes mid-generation?
A: FastAPI handles exceptions and returns appropriate HTTP error responses. The YAML isn't modified 
   during generation — only during a POST /save. So the data is safe.

Q: Is there any rate limiting or queue for the web UI?
A: No — this is designed as a local, single-user tool. Multi-user support is a future scope item.
""")
    return slide


def build_slide_10_folder(prs):
    """Slide 10: Folder Structure"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY_BG)
    add_slide_header(slide, "Project Structure", number=10)

    # Left: tree visual
    add_card(slide, Inches(0.4), Inches(1.1), Inches(6.0), Inches(6.15),
             fill=NAVY_CARD, radius=18000)

    tree_lines = [
        ("resume-engine/",                       Pt(13), ACCENT_CYAN, False, 0),
        ("├── compile_resume.py",                Pt(11.5), OFF_WHITE, False, 1),
        ("├── filter_engine.py",                 Pt(11.5), OFF_WHITE, False, 1),
        ("├── pdf_renderer.py",                  Pt(11.5), OFF_WHITE, False, 1),
        ("├── git_utils.py",                     Pt(11.5), OFF_WHITE, False, 1),
        ("├── run_web.py",                       Pt(11.5), OFF_WHITE, False, 1),
        ("├── requirements.txt",                 Pt(11.5), MUTED, False, 1),
        ("├── README.md",                        Pt(11.5), MUTED, False, 1),
        ("├── data/",                            Pt(11.5), ACCENT_BLUE, True, 1),
        ("│   ├── resume.example.yaml",          Pt(10.5), OFF_WHITE, False, 2),
        ("│   └── resume.yaml  (gitignored)",    Pt(10.5), MUTED, True, 2),
        ("├── templates/",                       Pt(11.5), ACCENT_BLUE, True, 1),
        ("│   └── resume_template.html",         Pt(10.5), MUTED, False, 2),
        ("├── tests/",                           Pt(11.5), _rgb(0x22,0xC5,0x5E), True, 1),
        ("│   ├── conftest.py",                  Pt(10.5), OFF_WHITE, False, 2),
        ("│   └── test_filter_engine.py",        Pt(10.5), OFF_WHITE, False, 2),
        ("├── web/",                             Pt(11.5), ACCENT_CYAN, True, 1),
        ("│   ├── app.py",                       Pt(10.5), OFF_WHITE, False, 2),
        ("│   ├── models.py",                    Pt(10.5), OFF_WHITE, False, 2),
        ("│   ├── data_store.py",               Pt(10.5), OFF_WHITE, False, 2),
        ("│   ├── templates/",                   Pt(10.5), ACCENT_CYAN, True, 2),
        ("│   └── static/",                      Pt(10.5), ACCENT_CYAN, True, 2),
        ("└── .github/workflows/ci.yml",         Pt(11.5), _rgb(0xFB,0xD5,0x68), False, 1),
    ]

    tree_y = Inches(1.22)
    for (line, sz, col, bold, _) in tree_lines:
        add_textbox(slide, line,
                    Inches(0.6), tree_y, Inches(5.6), Inches(0.26),
                    font_name="Aptos", font_size=sz, color=col, bold=bold)
        tree_y += Inches(0.248)

    # Right: annotations
    annotations = [
        ("CLI Entry Point",
         "Argparse-based CLI. Flags: --target, --output, --push, --preview, --list-targets, --targets-report.",
         Inches(6.6), Inches(1.15), ACCENT_BLUE),
        ("Core Engine",
         "filter_engine.py — pure functions, unit tested.\npdf_renderer.py — ReportLab canvas, multi-page overflow.\ngit_utils.py — auto-commit/push, gitignore-aware.",
         Inches(6.6), Inches(2.05), ACCENT_CYAN),
        ("Data Files",
         "resume.example.yaml — fictional demo data, committed.\nresume.yaml — your real data, gitignored (never committed).",
         Inches(6.6), Inches(3.15), ACCENT_BLUE),
        ("HTML Template (Dormant)",
         "Jinja2/WeasyPrint renderer — exists but not yet wired up. Listed in roadmap.",
         Inches(6.6), Inches(3.95), AMBER),
        ("Test Suite",
         "~40 pytest unit tests. conftest.py holds shared fixtures.",
         Inches(6.6), Inches(4.6), GREEN_OK),
        ("Web Module",
         "FastAPI routes (app.py) · Pydantic models · YAML data store · Jinja2 templates · CSS dark theme · vanilla JS.",
         Inches(6.6), Inches(5.15), ACCENT_CYAN),
        ("CI/CD",
         "GitHub Actions: runs pytest + smoke-compiles backend/data-science/frontend on every push.",
         Inches(6.6), Inches(6.05), _rgb(0xFB,0xD5,0x68)),
    ]

    for (title, desc, ax, ay, col) in annotations:
        add_shape(slide, ax - Inches(0.1), ay + Inches(0.1), Inches(0.04), Inches(0.5), fill_color=col)
        add_textbox(slide, title,
                    ax, ay, Inches(6.2), Inches(0.28),
                    font_name="Aptos Display", font_size=Pt(11.5), bold=True, color=col)
        add_textbox(slide, desc,
                    ax, ay + Inches(0.3), Inches(6.2), Inches(0.55),
                    font_name="Aptos", font_size=Pt(10), color=OFF_WHITE, wrap=True)

    add_speaker_notes(slide, """SPEAKER NOTES — Slide 10: Project Structure

WHAT TO SAY:
Let me walk through the folder structure. I've cleaned it up to show only the meaningful files — 
removing venv, __pycache__, and .pytest_cache.

ROOT LEVEL — Core modules:
- compile_resume.py: The CLI with 6 flags
- filter_engine.py: The heart of the engine — pure functions, fully tested
- pdf_renderer.py: ReportLab-based PDF drawing
- git_utils.py: Version control automation
- run_web.py: Simple script to start the FastAPI web server

DATA/ — Two YAML files:
- resume.example.yaml: Fictional person "Jordan Rivera" — committed to the repo, safe for CI
- resume.yaml: Your real data — gitignored, never committed (important for privacy)

TEMPLATES/ — One dormant HTML file:
- resume_template.html: A Jinja2 template that exists but isn't wired to the engine yet. 
  It's a roadmap item for a future WeasyPrint/HTML rendering path.

TESTS/ — Full pytest suite:
- conftest.py: Shared fixtures (sample_data, tmp_yaml)
- test_filter_engine.py: ~40 tests covering all filter_engine functions

WEB/ — The FastAPI web application:
- app.py: Five routes. models.py: Pydantic schemas. data_store.py: YAML I/O.
- templates/ and static/: Jinja2 HTML, dark-theme CSS, vanilla JavaScript.

.github/workflows/ci.yml — GitHub Actions that runs on every push.

NOTE: There's a stray resume_data.yaml at the project root from an earlier iteration — 
it's not read by any current module and can be ignored.

EXPECTED FACULTY QUESTIONS:
Q: Why is the actual resume YAML gitignored?
A: Personal data protection. data/resume.yaml contains your real contact info, employer details, 
   and career history — you don't want that accidentally pushed to a public repository. 
   The fictional example ships instead.
""")
    return slide


def build_slide_11_demo(prs):
    """Slide 11: Demonstration (placeholder frames with real data reference)"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY_BG)
    add_slide_header(slide, "Live Demonstration", number=11)

    # Subtitle
    add_textbox(slide, "Replace placeholder frames below with actual screenshots before submission.",
                Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.28),
                font_name="Aptos", font_size=Pt(10.5), italic=True, color=MUTED)

    screenshot_frames = [
        ("Edit Resume Page", "http://127.0.0.1:8000/", "Open the web UI and screenshot the form editor.", ACCENT_BLUE),
        ("Generate PDF Page", "http://127.0.0.1:8000/generate", "Pick 'backend' target and screenshot the generate page.", ACCENT_CYAN),
        ("Generated PDF", "/tmp/resume_backend_demo.pdf", "Screenshot page 1 of the generated resume PDF.", GREEN_OK),
        ("--targets-report CLI", "python compile_resume.py\n--targets-report", "Screenshot the colored terminal table.", AMBER),
    ]

    w_card = Inches(3.0)
    h_card = Inches(4.2)
    gap = Inches(0.25)
    start_x = Inches(0.4)
    card_y = Inches(1.35)

    for i, (title, url, instruction, col) in enumerate(screenshot_frames):
        cx = start_x + i * (w_card + gap)
        add_card(slide, cx, card_y, w_card, h_card,
                 fill=NAVY_CARD, border=col, radius=18000)
        add_shape(slide, cx, card_y, w_card, Inches(0.055), fill_color=col)

        # Screenshot placeholder area
        add_card(slide, cx + Inches(0.12), card_y + Inches(0.12),
                 w_card - Inches(0.24), Inches(2.6),
                 fill=_rgb(0x0A,0x0F,0x18), border=_rgb(0x25,0x35,0x50), radius=12000)
        add_textbox(slide, "[ Screenshot Placeholder ]",
                    cx + Inches(0.12), card_y + Inches(0.12),
                    w_card - Inches(0.24), Inches(2.6),
                    font_name="Aptos", font_size=Pt(10), italic=True,
                    color=_rgb(0x30,0x45,0x65), align=PP_ALIGN.CENTER)

        # Title
        add_textbox(slide, title,
                    cx + Inches(0.12), card_y + Inches(2.78),
                    w_card - Inches(0.15), Inches(0.4),
                    font_name="Aptos Display", font_size=Pt(12), bold=True, color=WHITE)

        # URL/command
        add_textbox(slide, url,
                    cx + Inches(0.12), card_y + Inches(3.2),
                    w_card - Inches(0.15), Inches(0.45),
                    font_name="Aptos", font_size=Pt(9), italic=True, color=col)

        # Instruction
        add_textbox(slide, instruction,
                    cx + Inches(0.12), card_y + Inches(3.68),
                    w_card - Inches(0.15), Inches(0.45),
                    font_name="Aptos", font_size=Pt(8.5), color=MUTED, wrap=True)

    # How to capture instructions
    add_card(slide, Inches(0.4), Inches(5.72), Inches(12.5), Inches(1.6),
             fill=_rgb(0x14,0x1E,0x30), radius=15000)
    add_textbox(slide, "📸  How to Capture Screenshots Before Submission:",
                Inches(0.6), Inches(5.77), Inches(12.0), Inches(0.32),
                font_name="Aptos Display", font_size=Pt(12), bold=True, color=ACCENT_CYAN)
    steps = "1. cd resume-engine  &&  source venv/bin/activate  &&  python run_web.py\n" \
            "2. Browser → http://127.0.0.1:8000  →  screenshot Edit page\n" \
            "3. Browser → http://127.0.0.1:8000/generate  →  screenshot Generate page\n" \
            "4. Click Generate (target: backend)  →  open downloaded PDF  →  screenshot page 1\n" \
            "5. Terminal: python compile_resume.py --targets-report  →  screenshot colored table"
    add_textbox(slide, steps,
                Inches(0.6), Inches(6.12), Inches(12.0), Inches(1.15),
                font_name="Aptos", font_size=Pt(10), color=OFF_WHITE, wrap=True)

    add_speaker_notes(slide, """SPEAKER NOTES — Slide 11: Live Demonstration

OPTION A — LIVE DEMO (preferred):
If you can connect to a computer running the project, do a live demo:
1. Open terminal → python compile_resume.py --targets-report (show colored table)
2. Open browser → http://127.0.0.1:8000 (show Edit Resume page)
3. Navigate to /generate, select 'backend', click Generate (download PDF)
4. Open the downloaded PDF and show it on screen

OPTION B — SCREENSHOT DEMO (fallback):
Replace the 4 placeholder frames on this slide with actual screenshots:
1. Edit Resume page (the web form with pre-filled career data)
2. Generate PDF page (target dropdown, git commit checkbox)
3. Generated PDF page 1 (Jordan Rivera's backend resume)
4. Terminal showing --targets-report output (the colored table)

TALKING POINTS DURING DEMO:
• "Notice the form is pre-filled — it reads from the YAML file automatically."
• "I'm selecting 'backend' as the target. This will filter out any frontend or data-science items."
• "The PDF generates in under a second — no browser, no external tools."
• "The targets-report shows exactly how many items each role would get — useful for checking 
   your YAML data coverage."

EXPECTED FACULTY QUESTIONS:
Q: Can we see the YAML file?
A: Yes — open data/resume.example.yaml and show a few entries with their tags.
Q: How fast is PDF generation?
A: Under 1 second for a typical resume. ReportLab's canvas API is very fast.
""")
    return slide


def build_slide_12_challenges(prs):
    """Slide 12: Challenges & Solutions"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY_BG)
    add_slide_header(slide, "Engineering Challenges", number=12)

    challenges = [
        ("Multi-Page PDF Overflow",
         "Problem: Four separate try/except blocks for overflow — brittle and repetitive.",
         "Solution: Unified _check_or_new_page() overflow gate — one function handles page breaks throughout the renderer.",
         "📄"),
        ("ReportLab Word-Wrap & Hanging Indents",
         "Problem: ReportLab's canvas API has no built-in rich text flow — bullet wrapping and indentation required manual calculation.",
         "Solution: Explicit y-coordinate tracking and manual text-width measurement for correct hanging indent rendering.",
         "📐"),
        (".gitignore-Aware Auto-Commit",
         "Problem: GitPython reports an already-tracked file as 'not ignored' even if it matches a .gitignore rule — a quiet bug.",
         "Solution: Explicit safeguard in git_utils.py — checks tracking status AND .gitignore independently before committing.",
         "🔗"),
        ("Two Different 'Empty' Behaviors",
         "Problem: A job with zero matching bullets should be dropped entirely; a degree entry with zero matching coursework should be kept.",
         "Solution: filter_experience() drops jobs with no bullets; filter_education() always keeps the degree entry, only filtering coursework within it.",
         "🔀"),
        ("CLI + Web UI Sync",
         "Problem: Risk of the CLI and web UI using different filter/render logic — any divergence is a silent bug.",
         "Solution: Both frontends explicitly import and call the same filter_engine.py and pdf_renderer.py — zero logic duplication.",
         "⚙️"),
        ("Protecting Personal Data",
         "Problem: data/resume.yaml contains real personal info — accidental git commit would expose it.",
         "Solution: data/resume.yaml is gitignored; a fictional example dataset (Jordan Rivera) ships for safe demo and CI use.",
         "🔒"),
    ]

    w = Inches(6.1)
    h = Inches(1.5)
    gap_x = Inches(0.5)
    gap_y = Inches(0.22)
    cols = [Inches(0.4), Inches(6.85)]
    rows = [Inches(1.1), Inches(2.85), Inches(4.6)]

    for i, (title, problem, solution, icon) in enumerate(challenges):
        col = i % 2
        row = i // 2
        cx = cols[col]
        cy = rows[row]

        add_card(slide, cx, cy, w, h, fill=NAVY_CARD, border=_rgb(0x25,0x3A,0x5E), radius=18000)
        add_shape(slide, cx, cy, w, Inches(0.05), fill_color=ACCENT_BLUE)

        # Icon + title
        add_textbox(slide, icon, cx + Inches(0.15), cy + Inches(0.1),
                    Inches(0.4), Inches(0.42), font_size=Pt(18))
        add_textbox(slide, title,
                    cx + Inches(0.58), cy + Inches(0.1), Inches(5.35), Inches(0.42),
                    font_name="Aptos Display", font_size=Pt(12.5), bold=True, color=WHITE)

        # Problem
        add_textbox(slide, f"⚠  {problem}",
                    cx + Inches(0.15), cy + Inches(0.58), Inches(5.8), Inches(0.4),
                    font_name="Aptos", font_size=Pt(9.5), color=AMBER, wrap=True)

        # Solution
        add_textbox(slide, f"✓  {solution}",
                    cx + Inches(0.15), cy + Inches(0.98), Inches(5.8), Inches(0.45),
                    font_name="Aptos", font_size=Pt(9.5), color=_rgb(0x22,0xC5,0x5E), wrap=True)

    add_speaker_notes(slide, """SPEAKER NOTES — Slide 12: Engineering Challenges

WHAT TO SAY:
Every non-trivial project has real engineering challenges. These six are genuine — pulled from 
code comments and commit history, not invented.

1. MULTI-PAGE OVERFLOW:
The first attempt had four separate try/except blocks checking for page overflow at different 
points in the renderer. This was fragile and repetitive. I refactored it into a single 
_check_or_new_page() function that centralizes all overflow logic. This is a classic DRY principle application.

2. REPORTLAB WORD-WRAP:
ReportLab's canvas draws text at coordinates — it doesn't have a "word wrap" mode for complex layouts. 
Getting bullet points to wrap correctly with a hanging indent (where the second line aligns under the 
text, not the bullet character) required tracking y-coordinates manually and measuring text widths.

3. GITIGNORE-AWARE COMMIT:
This was a subtle bug: GitPython's is_ignored() method returns False for a file that's already 
tracked, even if it matches a .gitignore rule. This meant the safeguard against committing gitignored 
files wasn't actually working in all cases. Added an explicit check for tracking status.

4. TWO 'EMPTY' BEHAVIORS:
A job entry and a degree entry have opposite intended behaviors when all their child items are filtered out. 
A job with zero bullets is entirely irrelevant — drop it. A degree entry is always relevant — keep it 
even with no coursework. These needed separate handling in filter_experience() vs filter_education().

5. CLI/WEB SYNC:
The architectural answer to "what if the two frontends drift apart?" — make it structurally impossible 
by having them both import from the same module.

6. PERSONAL DATA PROTECTION:
A practical DevOps concern. gitignoring the real data and shipping demo data is the correct solution.

EXPECTED FACULTY QUESTIONS:
Q: What debugging tools did you use?
A: Python's standard pdb debugger, plus the pytest test suite caught most regressions immediately.
""")
    return slide


def build_slide_13_scope(prs):
    """Slide 13: Future Scope"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY_BG)
    add_slide_header(slide, "Future Scope & Roadmap", number=13)

    # From README
    add_textbox(slide, "📋  From the Project's README Roadmap",
                Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.35),
                font_name="Aptos Display", font_size=Pt(15), bold=True, color=ACCENT_CYAN)

    roadmap = [
        ("🎨", "HTML/WeasyPrint Renderer",
         "Wire the dormant templates/resume_template.html with WeasyPrint for CSS-driven, browser-faithful PDF rendering. A --renderer html flag."),
        ("🎭", "Multi-Theme Support",
         "The web UI 'Theme' dropdown already exists (currently 'Coming soon'). Build a theme registry in pdf_renderer.py."),
        ("📐", "YAML Schema Validation (CLI)",
         "Add pre-render YAML validation using jsonschema or Pydantic to the compile path — currently only the web editor validates."),
        ("🖱️", "Interactive TUI",
         "A rich + prompt_toolkit terminal UI for browsing and selecting targets interactively — no need to remember tag names."),
        ("✉️", "Cover Letter Integration",
         "A parallel cover_letter.yaml + Jinja2 template using the same tag system for tailored cover letter generation."),
    ]

    for i, (icon, title, desc) in enumerate(roadmap):
        ry = Inches(1.55) + i * Inches(0.88)
        add_card(slide, Inches(0.4), ry, Inches(12.5), Inches(0.78),
                 fill=NAVY_CARD, border=ACCENT_BLUE, radius=15000)
        add_shape(slide, Inches(0.4), ry, Inches(0.06), Inches(0.78), fill_color=ACCENT_BLUE)
        add_textbox(slide, icon, Inches(0.55), ry, Inches(0.45), Inches(0.78),
                    font_size=Pt(20), color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, title,
                    Inches(1.05), ry + Inches(0.06), Inches(3.0), Inches(0.38),
                    font_name="Aptos Display", font_size=Pt(13), bold=True, color=WHITE)
        add_textbox(slide, desc,
                    Inches(4.1), ry + Inches(0.1), Inches(8.6), Inches(0.58),
                    font_name="Aptos", font_size=Pt(11.5), color=OFF_WHITE, wrap=True)

    # Stretch ideas
    add_textbox(slide, "💡  Reasonable Extension Ideas  (beyond the README)",
                Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.32),
                font_name="Aptos Display", font_size=Pt(13), bold=True, color=AMBER)
    add_card(slide, Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.82),
             fill=_rgb(0x28,0x22,0x10), radius=15000)
    stretch = "☁  Cloud hosting + authentication for multi-user accounts  ·  🤖  ATS-score estimation  ·  📌  Job-description-driven tag suggestions  ·  📊  Analytics dashboard"
    add_textbox(slide, stretch,
                Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.72),
                font_name="Aptos", font_size=Pt(11.5), color=AMBER, wrap=True)

    add_speaker_notes(slide, """SPEAKER NOTES — Slide 13: Future Scope

WHAT TO SAY:
The roadmap items on this slide come DIRECTLY from the README — these aren't vague wishes, 
they're planned next steps with concrete implementation approaches.

1. HTML/WeasyPrint Renderer:
The Jinja2 HTML template (templates/resume_template.html) already EXISTS in the repo — it just 
isn't connected to the engine yet. Future work would add a --renderer html flag. The advantage 
is CSS-based design control that ReportLab's coordinate API can't easily replicate.

2. Multi-Theme Support:
The web UI already has a "Theme" dropdown — it just says "Coming soon." The renderer currently 
has one hardcoded set of design tokens (fonts, colors, margins). Building a theme registry 
is the next logical step.

3. YAML Schema Validation (CLI):
Currently, Pydantic validation only runs in the web editor's save path. The CLI compile path 
accepts any YAML structure and handles errors defensively with .get() and defaults. Adding 
pre-render validation would give better error messages for malformed YAML files.

4. Interactive TUI:
Right now you need to know the tag names to use the CLI. A rich terminal UI would let you 
browse available targets interactively.

5. Cover Letter Integration:
The tag system generalizes naturally to cover letters — you'd maintain a cover_letter.yaml with 
role-tagged paragraphs.

[For the stretch ideas, frame them clearly as your own extension ideas, not from the README:]
These are ideas I think would make logical next steps beyond the roadmap: cloud hosting with auth, 
ATS scoring, and job-description parsing to suggest tags.

EXPECTED FACULTY QUESTIONS:
Q: Which of these would you implement next?
A: The HTML/WeasyPrint renderer — because the template already exists and it would dramatically 
   expand the design possibilities. The multi-theme support would follow naturally.
""")
    return slide


def build_slide_14_learnings(prs):
    """Slide 14: Key Learnings"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY_BG)
    add_slide_header(slide, "Key Learnings", number=14)

    learnings = [
        ("Separation of Concerns",
         "One core engine reused by two different frontends (CLI + web). Changes in filter_engine.py benefit both simultaneously.",
         "🏗️", ACCENT_BLUE),
        ("Pure, Testable Functions",
         "filter_engine.py has no side effects — all functions take data in, return data out. This made ~40 unit tests straightforward to write.",
         "🧪", _rgb(0x22,0xC5,0x5E)),
        ("Defensive Design",
         "Safe fallbacks throughout: demo data when real YAML is missing, gitignore-aware commit skip, .get() with defaults instead of hard crashes.",
         "🛡️", AMBER),
        ("REST API Design with FastAPI",
         "GET vs POST semantics, Pydantic for request validation, streaming responses for PDF download, structured error handling.",
         "⚡", ACCENT_CYAN),
        ("Schema Validation with Pydantic",
         "Validating structured user input at the API boundary — catching malformed data before it reaches the business logic layer.",
         "✅", _rgb(0x22,0xC5,0x5E)),
        ("Git Automation with GitPython",
         "Programmatic commits, push, and .gitignore inspection. Understanding Git internals (tracking status vs. ignore rules).",
         "🔗", _rgb(0xF4,0x73,0x73)),
        ("CI/CD with GitHub Actions",
         "Automated test runs and smoke-compilation on every push. The CI is what ensures the demo data covers all supported targets.",
         "🚀", _rgb(0xFB,0xD5,0x68)),
        ("Pragmatic Technology Choices",
         "ReportLab over HTML rendering (no browser dependency). FastAPI over Flask (async, built-in validation). Plain YAML over a database.",
         "⚖️", ACCENT_BLUE),
    ]

    w = Inches(6.1)
    h = Inches(1.35)
    cols = [Inches(0.4), Inches(6.85)]
    rows = [Inches(1.1), Inches(2.6), Inches(4.1), Inches(5.6)]

    for i, (title, desc, icon, col) in enumerate(learnings):
        ci = i % 2
        ri = i // 2
        cx = cols[ci]
        cy = rows[ri]

        add_card(slide, cx, cy, w, h, fill=NAVY_CARD, border=col, radius=18000)
        add_shape(slide, cx, cy, Inches(0.06), h, fill_color=col)

        add_textbox(slide, icon, cx + Inches(0.15), cy + Inches(0.1),
                    Inches(0.42), Inches(0.45), font_size=Pt(20))
        add_textbox(slide, title,
                    cx + Inches(0.62), cy + Inches(0.1), Inches(5.3), Inches(0.42),
                    font_name="Aptos Display", font_size=Pt(13), bold=True, color=WHITE)
        add_textbox(slide, desc,
                    cx + Inches(0.15), cy + Inches(0.6), Inches(5.8), Inches(0.68),
                    font_name="Aptos", font_size=Pt(10.5), color=OFF_WHITE, wrap=True)

    add_speaker_notes(slide, """SPEAKER NOTES — Slide 14: Key Learnings

WHAT TO SAY:
This project taught me several software engineering principles that I now understand from first-hand 
experience building a real system.

Walk through 3-4 key ones with elaboration:

SEPARATION OF CONCERNS:
This is the most valuable lesson. The decision to route both frontends through the same core modules 
meant that every improvement to filter_engine.py automatically benefited both the CLI and the web UI. 
This isn't just good theory — it's what kept the project maintainable as it grew.

PURE, TESTABLE FUNCTIONS:
filter_engine.py has NO side effects — no file I/O, no global state, no randomness. Every function 
takes some input and returns some output. This is EXACTLY why I could write 40 meaningful unit tests 
for it. If the functions had side effects, testing would have required complex mocking.

DEFENSIVE DESIGN:
Real software needs graceful degradation. Three examples from this project:
- No real YAML? Fall back to demo data, warn the user, keep working.
- Not in a git repo? git_utils gracefully skips the commit instead of crashing.
- Missing tags key in YAML item? filter_engine handles it silently with .get() and defaults.

PRAGMATIC TECH CHOICES:
Every technology choice was made for a concrete reason. ReportLab because it has no browser 
dependency. FastAPI because it has first-class Pydantic integration. Plain YAML because 
the data model is one user's career history — a database would be over-engineering.

EXPECTED FACULTY QUESTIONS:
Q: What was the most difficult concept to grasp?
A: The .gitignore behavior with GitPython — understanding that tracking status and ignore rules 
   are independent concepts in Git, and that a tracked file isn't "ignored" even if it matches 
   a .gitignore pattern.
""")
    return slide


def build_slide_15_testing(prs):
    """Slide 15: Test Coverage"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY_BG)
    add_slide_header(slide, "Testing & Engineering Practices", number=15)

    # Coverage matrix
    add_textbox(slide, "Test Coverage Reality Check  (engineering maturity requires honesty)",
                Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.35),
                font_name="Aptos Display", font_size=Pt(14), bold=True, color=WHITE)

    # Covered
    add_card(slide, Inches(0.4), Inches(1.52), Inches(6.0), Inches(4.1),
             fill=_rgb(0x0E,0x20,0x16), border=GREEN_OK, radius=18000)
    add_shape(slide, Inches(0.4), Inches(1.52), Inches(6.0), Inches(0.06), fill_color=GREEN_OK)
    add_textbox(slide, "✅  Fully Covered by ~40 Unit Tests",
                Inches(0.55), Inches(1.6), Inches(5.7), Inches(0.38),
                font_name="Aptos Display", font_size=Pt(14), bold=True, color=GREEN_OK)

    covered = [
        "filter_list() / _matches()  — tag inclusion & 'general' fallback",
        "rank_by_relevance()  — relevance scoring and priority: field",
        "apply_budget()  — trim, no-op, and exact-budget cases",
        "filter_experience()  — regression: zero-bullet jobs dropped",
        "filter_education()  — degrees always kept, coursework filtered",
        "filter_skills()  — empty categories omitted",
        "collect_all_tags()  — 'general' excluded, sorted, deduplicated",
        "build_resume_data()  — full integration test via tmp YAML fixture",
        "Edge cases  — empty lists, missing 'tags' key, unknown target",
    ]
    for j, item in enumerate(covered):
        add_textbox(slide, f"  ▸  {item}",
                    Inches(0.55), Inches(2.1) + j*Inches(0.35), Inches(5.7), Inches(0.32),
                    font_name="Aptos", font_size=Pt(10.5), color=OFF_WHITE)

    # Not yet covered
    add_card(slide, Inches(6.7), Inches(1.52), Inches(6.0), Inches(2.3),
             fill=_rgb(0x28,0x18,0x10), border=AMBER, radius=18000)
    add_shape(slide, Inches(6.7), Inches(1.52), Inches(6.0), Inches(0.06), fill_color=AMBER)
    add_textbox(slide, "⚠  Not Yet Unit Tested  (Future Work)",
                Inches(6.85), Inches(1.6), Inches(5.7), Inches(0.38),
                font_name="Aptos Display", font_size=Pt(14), bold=True, color=AMBER)

    not_covered = [
        "pdf_renderer.py  — ReportLab canvas output is hard to assert programmatically",
        "git_utils.py  — requires a live git repo fixture",
        "web/app.py  — FastAPI routes / HTTP integration tests",
    ]
    for j, item in enumerate(not_covered):
        add_textbox(slide, f"  ▸  {item}",
                    Inches(6.85), Inches(2.1) + j*Inches(0.4), Inches(5.7), Inches(0.35),
                    font_name="Aptos", font_size=Pt(10.5), color=OFF_WHITE)

    # CI/CD
    add_card(slide, Inches(6.7), Inches(3.98), Inches(6.0), Inches(1.65),
             fill=_rgb(0x14,0x20,0x30), border=_rgb(0xFB,0xD5,0x68), radius=18000)
    add_shape(slide, Inches(6.7), Inches(3.98), Inches(6.0), Inches(0.06),
              fill_color=_rgb(0xFB,0xD5,0x68))
    add_textbox(slide, "🚀  GitHub Actions CI Pipeline",
                Inches(6.85), Inches(4.06), Inches(5.7), Inches(0.38),
                font_name="Aptos Display", font_size=Pt(14), bold=True,
                color=_rgb(0xFB,0xD5,0x68))
    ci_steps = [
        "1. Setup: checkout · Python 3.12 · pip install -r requirements.txt",
        "2. Tests: pytest tests/ -v  (full ~40-test suite)",
        "3. Smoke: compile --target backend / data-science / frontend",
        "4. Verify: --list-targets · --targets-report  (no exceptions)",
    ]
    for j, s in enumerate(ci_steps):
        add_textbox(slide, s,
                    Inches(6.85), Inches(4.52) + j*Inches(0.28), Inches(5.7), Inches(0.26),
                    font_name="Aptos", font_size=Pt(10.5), color=OFF_WHITE)

    # Conftest callout
    add_card(slide, Inches(0.4), Inches(5.72), Inches(12.5), Inches(0.65),
             fill=NAVY_CARD, radius=12000)
    add_textbox(slide,
                "📌  Test fixtures are shared via tests/conftest.py — sample_data (in-memory dict) and tmp_yaml (writes to tmp_path). Adding a new target only requires updating conftest.py.",
                Inches(0.6), Inches(5.72), Inches(12.1), Inches(0.65),
                font_name="Aptos", font_size=Pt(11), color=OFF_WHITE, wrap=True)

    # Future testing roadmap
    add_card(slide, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.78),
             fill=_rgb(0x1A,0x1A,0x30), radius=12000)
    add_textbox(slide,
                "🔮  Planned: pytest-httpx for FastAPI integration tests · PDF content assertions via PyMuPDF · temporary git repo fixtures for git_utils",
                Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.78),
                font_name="Aptos", font_size=Pt(11), italic=True, color=MUTED, wrap=True)

    add_speaker_notes(slide, """SPEAKER NOTES — Slide 15: Testing & Engineering Practices

WHAT TO SAY:
This slide is about engineering maturity and intellectual honesty. I could have just said 
"the project has unit tests" — but I think it's more valuable to be precise about what IS and 
ISN'T tested.

WHAT IS COVERED: ~40 tests cover filter_engine.py comprehensively. Every function has multiple 
test cases, including edge cases. The conftest.py provides shared fixtures so tests are 
concise and focused. I have regression tests — for example, a specific test that locks in 
the behavior that "a job with zero matching bullets is dropped entirely."

WHAT IS NOT COVERED: Three modules currently have no unit tests:
1. pdf_renderer.py — difficult to test programmatically because the output is a binary PDF. 
   You'd need a library like PyMuPDF to inspect the rendered content.
2. git_utils.py — requires a live git repo. You'd need pytest fixtures that create temporary 
   git repositories.
3. web/app.py — requires HTTP-level integration testing (pytest-httpx or similar).

WHY I'M NAMING THIS EXPLICITLY:
Knowing what you HAVEN'T tested is part of engineering responsibility. Pretending everything 
is tested when it isn't creates false confidence.

CI/CD: The GitHub Actions workflow runs on every push and has three stages:
1. Full test suite
2. Smoke-compiles three targets (backend, data-science, frontend)
3. Runs --list-targets and --targets-report to verify CLI flags work

EXPECTED FACULTY QUESTIONS:
Q: What is a "regression test"?
A: A test written specifically to lock in a behavior that was once a bug or a design decision. 
   If someone accidentally changes the code so that zero-bullet jobs are no longer dropped, 
   the regression test will catch it immediately.

Q: How do you test the PDF output without rendering it?
A: Currently, it's not tested — that's acknowledged as a gap. The smoke test in CI does 
   compile the PDF and would fail if the renderer throws an exception, but it doesn't 
   assert anything about the PDF's content or layout.
""")
    return slide


def build_slide_16_conclusion(prs):
    """Slide 16: Conclusion"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, NAVY_BG)

    # Top accent bar
    add_shape(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=ACCENT_BLUE)
    add_shape(slide, 0, 0, Inches(0.06), SLIDE_H, fill_color=ACCENT_BLUE)

    # Decorative circles (same as title slide)
    add_shape(slide, Inches(10.8), Inches(-0.5), Inches(3.2), Inches(3.2),
              fill_color=_rgb(0x1E,0x3A,0x5F), radius=100000)
    add_shape(slide, Inches(11.2), Inches(-0.1), Inches(2.4), Inches(2.4),
              fill_color=_rgb(0x25,0x63,0xEB), radius=100000)
    add_shape(slide, Inches(11.6), Inches(0.3), Inches(1.6), Inches(1.6),
              fill_color=_rgb(0x38,0xBD,0xF8), radius=100000)

    # 4 conclusion lines
    items = [
        ("🎯 Problem",
         "Resume customization is repetitive, error-prone, and produces inconsistent versions of your career history."),
        ("💡 Solution",
         "Resume Engine: one YAML source of truth, tag-based filtering, and automatic PDF generation via CLI or web UI."),
        ("⚡ Impact",
         "Role-specific, one-page PDF resumes in under a second — with git version history, CI validation, and a Pydantic-validated web editor."),
        ("📚 Learning",
         "Separation of concerns · pure functions · defensive design · REST APIs · Git automation · CI/CD — real engineering principles, applied."),
    ]

    for i, (label, text) in enumerate(items):
        y = Inches(1.5) + i * Inches(1.2)
        add_card(slide, Inches(0.5), y, Inches(10.5), Inches(1.0),
                 fill=NAVY_CARD, border=ACCENT_BLUE, radius=18000)
        add_shape(slide, Inches(0.5), y, Inches(0.07), Inches(1.0),
                  fill_color=[ACCENT_BLUE, ACCENT_CYAN, GREEN_OK, AMBER][i])
        add_textbox(slide, label,
                    Inches(0.7), y + Inches(0.08), Inches(2.2), Inches(0.38),
                    font_name="Aptos Display", font_size=Pt(14), bold=True,
                    color=[ACCENT_BLUE, ACCENT_CYAN, GREEN_OK, AMBER][i])
        add_textbox(slide, text,
                    Inches(2.95), y + Inches(0.1), Inches(8.0), Inches(0.78),
                    font_name="Aptos", font_size=Pt(13), color=OFF_WHITE, wrap=True)

    # Thank you
    add_textbox(slide, "Thank You",
                Inches(0.5), Inches(6.1), Inches(10.5), Inches(0.75),
                font_name="Aptos Display", font_size=Pt(40), bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    # Questions tagline
    add_textbox(slide, "Questions? I'd love to go deeper on any part of the implementation.",
                Inches(0.5), Inches(6.82), Inches(12.0), Inches(0.35),
                font_name="Aptos", font_size=Pt(13), italic=True,
                color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

    # Bottom bar
    add_shape(slide, 0, Inches(7.3), SLIDE_W, Inches(0.2), fill_color=NAVY_CARD)
    add_textbox(slide, "Resume Engine  ·  Summer Training Project  ·  [Author Name]  ·  [College]",
                Inches(0.55), Inches(7.3), Inches(11), Inches(0.2),
                font_size=Pt(9), color=MUTED)

    add_speaker_notes(slide, """SPEAKER NOTES — Slide 16: Conclusion

WHAT TO SAY:
Let me close with a summary of what this project demonstrates:

PROBLEM: Resume customization is tedious, error-prone, and produces inconsistent versions. 
We needed a systematic solution.

SOLUTION: Resume Engine is a Git-driven, tag-based compilation system. One YAML file, 
unlimited tailored resumes. CLI for developers, web UI for visual editing.

IMPACT: Role-specific PDFs in under a second. Git version history of every resume. 
GitHub Actions CI that validates the engine on every push. A Pydantic-validated web editor 
that prevents data corruption.

LEARNING: This project gave me hands-on experience with real software engineering principles: 
separation of concerns through the shared core, pure functions that enable unit testing, 
defensive design with graceful fallbacks, REST API design with FastAPI, Git automation with 
GitPython, and CI/CD with GitHub Actions.

CLOSING:
I'm happy to take any questions — whether about the architecture decisions, specific implementation 
details, the challenges I faced, or the roadmap. I can also do a live demonstration if there's 
interest.

EXPECTED FACULTY QUESTIONS:
Q: What grade/rating would you give this project?
A: [Let them assess! Don't answer this one yourself.]

Q: What would you do differently if you were starting over?
A: I would add the HTML/WeasyPrint renderer from the beginning rather than building it in parallel 
   with ReportLab. Starting with one renderer and planning the theme system from day one 
   would have made the architecture cleaner.

Q: Is this deployed anywhere?
A: Currently it's a local tool only — deployment is a roadmap item. Adding authentication for 
   multi-user use would be the prerequisite for cloud deployment.
""")
    return slide


# ─── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building Resume Engine Presentation...")

    builders = [
        (build_slide_01_title,        "Slide 01 — Title"),
        (build_slide_02_problem,      "Slide 02 — Problem Statement"),
        (build_slide_03_overview,     "Slide 03 — Project Overview"),
        (build_slide_04_features,     "Slide 04 — Key Features"),
        (build_slide_05_techstack,    "Slide 05 — Technology Stack"),
        (build_slide_06_architecture, "Slide 06 — Architecture"),
        (build_slide_07_modules,      "Slide 07 — Module Architecture"),
        (build_slide_08_dataflow,     "Slide 08 — Data Flow"),
        (build_slide_09_workflow,     "Slide 09 — Workflow"),
        (build_slide_10_folder,       "Slide 10 — Folder Structure"),
        (build_slide_11_demo,         "Slide 11 — Demonstration"),
        (build_slide_12_challenges,   "Slide 12 — Challenges"),
        (build_slide_13_scope,        "Slide 13 — Future Scope"),
        (build_slide_14_learnings,    "Slide 14 — Key Learnings"),
        (build_slide_15_testing,      "Slide 15 — Testing & CI"),
        (build_slide_16_conclusion,   "Slide 16 — Conclusion"),
    ]

    for fn, label in builders:
        print(f"  {label}...", end=" ", flush=True)
        try:
            fn(prs)
            print("✓")
        except Exception as e:
            print(f"✗  ERROR: {e}")
            import traceback
            traceback.print_exc()

    output_path = "Resume_Engine_Presentation.pptx"
    prs.save(output_path)
    print(f"\n✓ Saved: {output_path}")
    print(f"  Slides: {len(prs.slides)}")
    return output_path


if __name__ == "__main__":
    main()
